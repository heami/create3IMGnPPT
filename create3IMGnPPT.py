from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageOps, ImageFilter, ImageGrab
import requests
from bs4 import BeautifulSoup
import glob
import os
import sys
import re
import createImgPPT_config as cfg
# --- 엑셀 이미지 생성을 위해 추가된 라이브러리 ---
import glob
import time
import gc
import pythoncom
import win32com.client
from win32com.client.gencache import EnsureDispatch
from win32com.client import constants
import fitz  # PyMuPDF

# 스크립트 파일이 있는 디렉토리 경로를 미리 저장합니다.
script_dir = os.path.dirname(os.path.abspath(__file__))
template_path = os.path.join(script_dir, "ppt-template.pptx")

Image.MAX_IMAGE_PIXELS = None

def create_regular_images(directory):
    """
    주어진 디렉토리에서 원본 이미지를 찾아 -l, -m 2가지 스펙의 이미지로 변환합니다.
    """
    print('--- 이미지 생성 작업을 시작합니다. ---')

    # 디렉토리에서 원본 이미지 파일 목록을 가져옵니다.
    # .jpg 또는 .tif 파일이면서 파일 이름에 '-g' 또는 '-i'와 숫자가 포함된 이미지만 가져옵니다.
    image_files = [f for f in os.listdir(directory) if
                   (f.endswith('.jpg') or f.endswith('.tif')) and
                   re.search(r'-[gi]\d{3}', f) and
                   not f.endswith(('-m.jpg'))]

    # 가져온 이미지 파일 이름 리스트를 이름 오름차순으로 정렬합니다.
    image_files.sort()

    if not image_files:
        print("경고: 이미지 파일을 찾을 수 없습니다. (.jpg, .tif)")
        return False

    for file in image_files:
        try:
            print(f"    • {file}")
            full_path = os.path.join(directory, file)
            img = Image.open(full_path)

            imagename = file.replace('.tif', '.jpg')
            l_save_filename = imagename.replace('.jpg', '-l.jpg')
            m_save_filename = imagename.replace('.jpg', '-m.jpg')

            save_dir = os.path.join(directory, 'Images')  # 이미지를 Images 폴더에 저장한다고 가정

            if not os.path.exists(save_dir):
                os.makedirs(save_dir)

            # l-image, m-image만들기
            cfg.createlim(img, l_save_filename, save_dir)
            cfg.createmim(img, m_save_filename, save_dir)

        except Exception as e:
            print(f"오류: {file} 파일 변환 실패 - {e}")

    print('--- 이미지 생성 완료. ---')
    print('   ')
    return True

def create_excel_images(directory):
    """
    주어진 디렉토리에서 Excel 영역을 캡쳐하고 -l, -m 2가지 스펙의 이미지로 변환합니다.
    """
    print('--- Excel 이미지 생성 작업을 시작합니다. ---')

    excel_app = None
    try:
        # 백그라운드 스레드에서 COM 객체 사용을 위해 초기화
        pythoncom.CoInitialize()

        # --- 0. 엑셀 애플리케이션 실행. DispatchEx를 사용하여 다른 엑셀 창과 독립된 새 프로세스를 시작
        excel_app = EnsureDispatch('Excel.Application')
        excel_app.Visible = False  # Excel 창을 보이지 않게 실행
        excel_app.DisplayAlerts = False # Excel이 사용자에게 보여주는 모든 종류의 확인 창, 경고 메시지, 알림 팝업을 띄우지 않도록 설정

        # --- 1. 전처리 ---
        excelFilePathList = [file for file in glob.glob(os.path.join(directory, '*-i???.xlsx'))
                             if not os.path.basename(file).startswith('~$')]
        excelFilePathList.sort()

        images_dir = os.path.join(directory, 'Images')
        if not os.path.isdir(images_dir):
            os.mkdir(images_dir)

        # --- 2. 본작업 ---
        for excel_path in excelFilePathList:
            workbook = None
            sheet = None
            temp_sheet = None
            table_range = None
            pasted_picture = None
            original_gridline_state = None
            try:
                # 엑셀 파일 열기 (절대 경로 사용) 첫 번째 시트를 대상으로 작업
                abs_excel_path = os.path.abspath(excel_path)
                workbook = excel_app.Workbooks.Open(abs_excel_path, ReadOnly=True)
                sheet = workbook.Sheets(1)
                sheet.Activate()  # 작업을 위해 시트 활성화

                # --------------------- 눈금선 제어 로직 추가 ---------------------
                original_gridline_state = excel_app.ActiveWindow.DisplayGridlines
                excel_app.ActiveWindow.DisplayGridlines = False
                # ----------------------------------------------------------------

                # 데이터가 있는 마지막 행과 열 번호를 찾음.
                last_row = sheet.Cells(sheet.Rows.Count, 1).End(constants.xlUp).Row  # A열 기준 마지막 행
                bottom_border = sheet.Cells(sheet.Rows.Count, 1).End(constants.xlUp).Borders(constants.xlEdgeBottom)
                if bottom_border.LineStyle == constants.xlNone:
                    last_row = last_row - 1 # 마지막 행이 캡션인 경우
                else:
                    last_row

                last_col = 1  # 최소값 1로 초기화
                # 2. 마지막 열을 더 안정적으로 찾기: 모든 행을 검사하여 가장 넓은 열을 찾음
                for r in range(1, last_row + 1):
                    current_last_col = sheet.Cells(r, sheet.Columns.Count).End(constants.xlToLeft).Column
                    if current_last_col > last_col:
                        last_col = current_last_col  # 더 큰 값이 나오면 업데이트

                print(f"    • {os.path.basename(excel_path)} (마지막 행: {last_row}, 마지막 열: {last_col})")

                # 첫 행(제목)과 마지막 행(캡션)을 제외한 범위 객체를 생성.
                # 시작: 2행 1열 (A2) / 끝: 테이블 마지막 행, 마지막 열
                if last_row > 1 and last_col >= 1:  # 데이터가 타이틀 포함 2행 이상일 때만 의미가 있음
                    table_range = sheet.Range(sheet.Cells(2, 1), sheet.Cells(last_row, last_col))

                    # --------------------- '그림으로 붙여넣기'  ---------------------
                    # 1. 계산된 범위를 '그림'이 아닌 '데이터'로 복사
                    table_range.Copy()

                    # 2. 임시 시트를 마지막에 추가
                    temp_sheet = workbook.Sheets.Add(After=workbook.Sheets(workbook.Sheets.Count))
                    temp_sheet.Activate()

                    # 3. 임시 시트에 '그림'으로 붙여넣기
                    pasted_picture = temp_sheet.Pictures().Paste()

                    # 4. 붙여넣어진 '그림 개체'를 클립보드로 복사
                    pasted_picture.Copy()

                    # table_range.CopyPicture(Format=constants.xlBitmap)

                    # 5. 클립보드에서 RGBA 이미지 가져오기
                    img = ImageGrab.grabclipboard()

                    if img:
                        # print(f"    - 클립보드 이미지 정보: mode={img.mode}, size={img.size}")
                        # 1. 최종 이미지를 담을 흰색 배경의 '도화지'를 먼저 생성합니다.
                        final_image = Image.new("RGB", img.size, (255, 255, 255))

                        # 2. 클립보드에서 가져온 이미지의 모드에 따라 분기 처리
                        if img.mode == 'RGBA':
                            # 시나리오 1: 투명도(Alpha)가 있는 이미지 (개발 PC 환경)
                            # print("    - 처리 방식: RGBA 모드 확인, 알파 채널을 마스크로 합성")
                            final_image.paste(img, mask=img.getchannel('A'))
                        elif img.mode == 'RGB':
                            # 시나리오 2: 투명도 없는 RGB 이미지 (노트북 환경)
                            # print("    - 처리 방식: RGB 모드 확인, 검은색을 투명으로 처리하는 마스크 생성 후 합성")
                            # 검은색(0)에 가까운 픽셀은 투명(0)으로, 그 외는 불투명(255)하게 만드는 마스크 생성
                            bw_mask = img.convert('L').point(lambda p: 255 if p > 10 else 0, mode='1')
                            final_image.paste(img, mask=bw_mask)
                        else:
                            # 시나리오 3: 그 외 다른 모드의 이미지 (예: 'P' 모드)
                            # print(f"    - 처리 방식: {img.mode} 모드 확인, RGBA로 변환 후 마스크 합성")
                            rgba_equivalent = img.convert('RGBA')
                            final_image.paste(rgba_equivalent, mask=rgba_equivalent.getchannel('A'))
                        # --------------------------------------------------------------------

                        # 파일명 생성
                        base_name, _ = os.path.splitext(os.path.basename(excel_path))
                        l_imgName = f"{base_name}-l.jpg"
                        m_imgName = f"{base_name}-m.jpg"
                        l_img_path = os.path.join(images_dir, l_imgName)

                        # 변환된 RGB 이미지로 l-image, m-image 만들기
                        final_image.save(l_img_path, dpi=(300, 300), quality=100)
                        cfg.createmim(final_image, m_imgName, images_dir)
                    else:
                        print(f"    - 경고: 이미지를 클립보드에서 가져오지 못했습니다.")

            except Exception as e:
                print(f"  - 오류: {os.path.basename(excel_path)} 처리 중 에러 발생 - {e}")

            finally:
                # --- 작업이 끝나면(성공하든 실패하든) 눈금선을 원래 상태로 복구 ---
                if original_gridline_state is not None:
                    excel_app.ActiveWindow.DisplayGridlines = original_gridline_state
                if temp_sheet:
                    temp_sheet.Delete()
                if workbook:
                    workbook.Close(SaveChanges=False)

        print('--- 이미지 생성 완료. ---\n')
        return True

    except KeyboardInterrupt:
        print("\n프로그램 실행이 중단되었습니다.")
    except Exception as e:
        print(f"프로그램 실행 중 심각한 오류 발생: {e}")
        return False
    finally:
        if excel_app:
            excel_app.DisplayAlerts = True
            excel_app.Quit()

        # COM 객체 참조를 명시적으로 해제하여 Python의 가비지 컬렉터가 객체를 정리하도록 돕습니다.
        pasted_picture = None
        table_range = None
        temp_sheet = None
        sheet = None
        workbook = None
        excel_app = None

        # 가비지 컬렉션을 수동으로 호출하여 메모리 정리를 유도
        gc.collect()

        # 백그라운드 스레드에서 COM 초기화 해제
        pythoncom.CoUninitialize()

def _deduplicate_lines(sorted_lines, tolerance=1.5):
    if not sorted_lines:
        return []
    result = []
    group = [sorted_lines[0]]
    for y in sorted_lines[1:]:
        if y - group[-1] <= tolerance:
            group.append(y)
        else:
            result.append(sum(group) / len(group))
            group = [y]
    result.append(sum(group) / len(group))
    return result

def _extract_horizontal_lines(page, page_width, min_width_ratio=0.3):
    min_line_width = page_width * min_width_ratio
    h_lines = []
    for drawing in page.get_drawings():
        for item in drawing.get("items", []):
            if item[0] == 'l':
                p1, p2 = item[1], item[2]
                if abs(p1.y - p2.y) <= 1.0 and abs(p2.x - p1.x) >= min_line_width:
                    h_lines.append(p1.y)
            elif item[0] == 're':
                rect = item[1]
                # 얇은 수평 사각형(선처럼 생긴 것)만 추출: 높이 ≤ 5pt
                if rect.width >= min_line_width and rect.height <= 5:
                    h_lines.append(rect.y0)
                    h_lines.append(rect.y1)
    return _deduplicate_lines(sorted(h_lines))

def _find_table_regions(page):
    """
    페이지에서 'Table N' 텍스트 위치를 기준으로 테이블 영역 (top_y, bottom_y) 목록 반환.
    - 각 테이블 제목 아래 첫 번째 수평선 = 테이블 상단
    - 다음 테이블 제목 또는 페이지 끝 이전 마지막 수평선 = 테이블 하단
    """
    page_width = page.rect.width
    page_height = page.rect.height

    # 1. 'Table N' 제목의 y 하단 좌표 수집
    # get_text("dict")로 줄 단위 좌표를 얻어 모든 줄을 검사
    title_bottoms = []
    for block in page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]:
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            line_text = "".join(span["text"] for span in line.get("spans", []))
            if re.match(r'Table\s*\d+', line_text.strip(), re.IGNORECASE):
                title_bottoms.append(line["bbox"][3])  # 해당 줄의 y1
    title_bottoms.sort()

    if not title_bottoms:
        return []

    # 2. 페이지 수평선 전체 추출
    all_h_lines = _extract_horizontal_lines(page, page_width)
    if len(all_h_lines) < 2:
        print(f"      - 진단: Table 제목 발견({len(title_bottoms)}개) but 수평선 부족({len(all_h_lines)}개)")
        return []

    # 3. 각 제목 아래 범위의 수평선으로 테이블 영역 계산
    regions = []
    boundaries = title_bottoms + [page_height]
    for i, title_y in enumerate(title_bottoms):
        next_boundary = boundaries[i + 1]
        table_lines = [y for y in all_h_lines if title_y <= y < next_boundary]
        if len(table_lines) >= 2:
            regions.append((table_lines[0], table_lines[-1]))

    return regions

def create_PDF_table_images(directory):
    """
    주어진 디렉토리에서 *-???.pdf 파일을 찾아 테이블 영역을 감지하고
    -l, -m 2가지 스펙의 이미지로 변환합니다.
    """
    print('--- PDF 테이블 이미지 생성 작업을 시작합니다. ---')

    pdf_file_list = sorted(glob.glob(os.path.join(directory, '*-???.pdf')))
    if not pdf_file_list:
        print("경고: *-???.pdf 패턴에 해당하는 파일을 찾을 수 없습니다.")
        return False

    images_dir = os.path.join(directory, 'Images')
    if not os.path.isdir(images_dir):
        os.mkdir(images_dir)

    zoom = 300 / 72

    for pdf_path in pdf_file_list:
        base_name, _ = os.path.splitext(os.path.basename(pdf_path))
        prefix = re.sub(r'-\d+$', '', base_name)
        print(f"    • {os.path.basename(pdf_path)}")
        doc = None
        try:
            doc = fitz.open(pdf_path)
            table_idx = 0

            for page in doc:
                table_groups = _find_table_regions(page)

                if not table_groups:
                    continue

                mat = fitz.Matrix(zoom, zoom)

                for top_y, bottom_y in table_groups:
                    table_idx += 1
                    clip_rect = fitz.Rect(0, top_y, page.rect.width, bottom_y)
                    pix = page.get_pixmap(matrix=mat, clip=clip_rect, alpha=False)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                    # 흰색 여백 제거
                    bbox = ImageOps.invert(img).getbbox()
                    img_cropped = img.crop(bbox) if bbox else img

                    # 상하좌우 10px 흰색 패딩 추가
                    w, h = img_cropped.size
                    final_image = Image.new("RGB", (w + 20, h + 20), (255, 255, 255))
                    final_image.paste(img_cropped, (10, 10))

                    # 파일명 생성 및 저장
                    img_base = f"{prefix}-i{table_idx:03d}"
                    l_img_path = os.path.join(images_dir, f"{img_base}-l.jpg")
                    final_image.save(l_img_path, dpi=(300, 300), quality=100)
                    cfg.createmim(final_image, f"{img_base}-m.jpg", images_dir)

            if table_idx == 0:
                print(f"      - 경고: {os.path.basename(pdf_path)}에서 테이블을 찾지 못했습니다.")

        except Exception as e:
            print(f"      - 오류: {os.path.basename(pdf_path)} 처리 중 에러 발생 - {e}")
        finally:
            if doc:
                doc.close()

    print('--- PDF 이미지 생성 완료. ---\n')
    return True

def dispatch_image_creation(directory):
    """엑셀 파일 존재 여부에 따라 적절한 이미지 생성 함수를 호출하는 분기점"""
    # print(f"'{os.path.basename(directory)}' 폴더에서 이미지 생성 작업을 시작합니다.")
    excel_pattern = os.path.join(directory, '*-i???.xlsx')
    pdf_pattern = os.path.join(directory, '*-???.pdf')
    excel_files = glob.glob(excel_pattern)

    if excel_files:
        # print("-> Excel 파일이 감지되었습니다. Excel 이미지 변환을 시작합니다.")
        return create_excel_images(directory)
    elif glob.glob(pdf_pattern):
        # print("-> Excel 파일이 감지되었습니다. Excel 이미지 변환을 시작합니다.")
        return create_PDF_table_images(directory)
    else:
        # print("-> Excel 파일이 없습니다. 일반 이미지 파일 처리를 시작합니다.")
        return create_regular_images(directory)





def create_ppts(directory, articleURLs, webURL, use_images_folder=False):
    """
    웹 스크래핑을 통해 논문 정보를 가져와 PPT를 생성합니다.
    use_images_folder: True일 경우 'Images' 폴더의 이미지를 사용합니다.
                       False일 경우 입력한 디렉토리의 이미지를 사용합니다.
    """
    print('--- PPT 생성 작업을 시작합니다. ---')

    # 사용할 이미지 폴더 경로를 결정합니다.
    if use_images_folder:
        image_dir = os.path.join(directory, 'Images')
    else:
        image_dir = directory

    # PPT를 저장할 'PPT' 폴더 경로
    ppt_dir = os.path.join(directory, 'PPT')

    if not os.path.isdir(ppt_dir):
        os.makedirs(ppt_dir)

    for url in articleURLs:
        try:
            response = requests.get(url)
            soup = BeautifulSoup(response.text, "html.parser")
            articlePage = \
            soup.select_one("#article-front-meta-left div").get_text().split('.')[1].strip().split(':')[1].split('-')[0]

            # FIG 정보 추출 (이전 코드와 동일한 로직)
            figOK = 0
            for label in soup.select("h5"):
                if label.get_text().find('Figures') >= 0:
                    figOK = 1

            if figOK == 0:
                # print(f"⚠️{url}  논문에 Figure가 없습니다.")
                continue

            figsSlider = soup.select_one(".figs-gray")
            getFigsURL = figsSlider.select(".carousel-inner a")

            figURLs = [webURL + fig['href'] for fig in getFigsURL]

            figLabels, figTitles, articleInfo, articleDois = [], [], [], []
            for fig_url in figURLs:
                # 웹 스크래핑으로 상세 정보 추출 (이전 코드와 동일한 로직)
                response = requests.get(fig_url)
                soup = BeautifulSoup(response.text, "html.parser")
                caption = soup.select('table')[2] if len(soup.select('table')) == 3 else soup.select('table')[1]

                figLabels.append(caption.select_one("td strong").get_text() + ' ')
                figTitles.append(caption.select_one("td span").get_text())
                info = caption.select_one(".smallText").get_text()
                articleInfo.append(info.split('.')[0] + info.split('.')[1] + '. ')

                doi = caption.select_one(".FigureTableDOI")
                articleDois.append(doi.get_text() if doi else '')

            # [수정] -m.jpg 파일을 제외하고 -l.jpg 파일 또는 .tif 파일만 사용하도록 파일 목록을 가져옵니다.
            filepaths, _ = cfg.getFileInfo(image_dir, articlePage)
            l_or_tif_filepaths = []

            # filepaths 리스트에 -l.jpg 파일이 있는지 확인합니다.
            has_l_jpg = any(p.endswith('-l.jpg') for p in filepaths)

            if has_l_jpg:
                # -l.jpg 파일이 하나라도 있으면, 해당 파일들만 사용합니다.
                l_or_tif_filepaths = [p for p in filepaths if p.endswith('-l.jpg')]
            else:
                # -l.jpg 파일이 없으면, .tif 파일들만 사용합니다.
                l_or_tif_filepaths = [p for p in filepaths if p.endswith('.tif')]

            # 가져온 이미지 파일 이름 리스트를 이름 오름차순으로 정렬합니다.
            l_or_tif_filepaths.sort()

            if not l_or_tif_filepaths:
                print(f"- {articlePage}에 해당하는 이미지 파일을 찾을 수 없습니다. PPT 생성을 건너뜁니다.")
                continue

            for file_path in l_or_tif_filepaths:
                match = re.search(r'-g(\d+)', os.path.basename(file_path))
                if match:
                    order = int(match.group(1)) - 1
                    try:
                        os.chdir(sys._MEIPASS)
                    except:
                        os.chdir(script_dir)

                    prs = Presentation(template_path)
                    img_info = cfg.px_to_inches(file_path)

                    hFactor = Inches(5) / img_info[1]
                    newHeight = Inches(5)
                    newWidth = img_info[0] * hFactor

                    if prs.slide_width < newWidth:
                        wFactor = Inches(9) / newWidth
                        newWidth = Inches(9)
                        newHeight = newHeight * wFactor

                    imgLeft = int((prs.slide_width - newWidth) / 2)
                    imgTop = int((prs.slide_height - newHeight) / 2)

                    prs.slides[0].shapes.add_picture(file_path, imgLeft, imgTop, height=newHeight)

                    cfg.inputInfo(figTitles[order], figLabels[order], articleInfo[order], articleDois[order], prs)

                    # PPT 저장 경로를 os.path.join으로 안전하게 구성
                    ppt_filename = os.path.basename(file_path).split('.')[0].replace('-l', '') + '.pptx'
                    save_path = os.path.join(ppt_dir, ppt_filename)
                    prs.save(save_path)
                    print(f"    • {ppt_filename}")

        except Exception as e:
            print(f"오류: {url} 처리 중 문제가 발생했습니다. - {e}")

    print('--- PPT 생성 완료. ---')

def createImgPPT_logic():
    # upcfg.check_for_updates()

    print('''
    ---------------------------------------------------------------
    * 작업 선택 *
    1. 이미지만 생성
    2. PPT만 생성
    3. 이미지와 PPT 모두 생성
    ---------------------------------------------------------------''')
    while True:
        mode = input("작업 모드를 선택해주세요 (1, 2, 3): ")
        if mode in ['1', '2', '3']:
            break
        print("잘못된 입력입니다. 1, 2, 3 중 하나를 입력해주세요.")

    directory = input("이미지가 있는 디렉토리 주소를 입력해주세요: ")

    if not os.path.isdir(directory):
        print("오류: 입력한 디렉토리가 존재하지 않습니다.")
        return

    os.chdir(directory)
    for d in ['PPT', 'Images']:
        if not os.path.isdir(os.path.join(directory, d)):
            os.mkdir(os.path.join(directory, d))

    if mode in ['2', '3']:
        URLToAnalyze = input("논문 page 또는 toc 주소를 입력해주세요: ")
        webURL = URLToAnalyze.split('/')[0] + '//' + URLToAnalyze.split('/')[2] + '/'
        articleURLs = cfg.decide_URL(URLToAnalyze, webURL)

    if mode in ['1', '3']:
        create_regular_images(directory)

    if mode in ['2', '3']:
        create_ppts(directory, articleURLs, webURL)

    print('-' * 100)
    input("작업이 완료되었습니다. 아무 키나 누르면 종료됩니다...")

# 프로그램의 시작점
if __name__ == "__main__":
    createImgPPT_logic()

