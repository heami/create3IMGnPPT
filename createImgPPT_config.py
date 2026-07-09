from pptx.util import Inches
from PIL import Image, ImageOps, ImageFilter
import requests
from bs4 import BeautifulSoup
import glob
import os
import re

### 함수선언
# toc, page 주소인 경우를 구분하여 url 리스트 반환
def decide_URL(url, webURL):
    articleURLs = []
    if 'asummary' in url:
        # toc 주소인 경우
        response = requests.get(url)
        soup = BeautifulSoup(response.text, "html.parser")

        for tag in soup.select(".ToC_item .ArticleInfor"):
            articleURLs.append(webURL + 'DOIx.php?id=' + tag.select_one("a").get_text()[16:])
    else:
        # page 주소인 경우
        articleURLs.append(url)
    # url 리스트를 반환
    return articleURLs

# 디렉토리 주소와 논문의 페이지를 받아 페이지와 동일한 fig의 이미지를 인식하여 파일주소와 파일이름 리스트를 반환
def getFileInfo(directory, articlePage):
    filepaths = []
    filenames = []

    for filename in os.listdir(directory):

        # 파일이 특정 패턴(articlePage)을 포함하고 있는지 확인
        if(re.search(r'-g\d{3}', filename) and (filename.endswith('.jpg') or filename.endswith('.tif'))):
            # 파일명에서 페이지 추출
            match = re.search(r'-([eS]?\d+)-[gi]\d+', filename)
            extracted_number = match.group(1)
            # 페이지가 일치하는지 확인
            if extracted_number == articlePage:
                # 이미지 파일의 전체 경로를 os.path.join으로 생성
                filepaths.append(os.path.join(directory, filename))
                filenames.append(filename)

    filepaths.sort()
    filenames.sort()

    # 'filepaths'는 절대 경로로 반환되어야 합니다.
    return filepaths, filenames

# image의 pixels 정보를 inches로 변환하여 그 값을 반환
def px_to_inches(path):
    im = Image.open(path)
    width = im.width / im.info['dpi'][0]  #해상도 145
    height = im.height / im.info['dpi'][1]  #해상도 145
    return (Inches(width), Inches(height))

# ppt에 해당하는 정보 삽입
def inputInfo(figTitle, figLabel, articleInfo, articleDois, prs):
    # 삽입될 글자수 제한
    # title 첫글자가 볼드나 이탤릭인 경우 \n\n로 시작함. 보정
    inputTitle = figTitle.replace('\n\n', '\n').split('\n')[1]
    if len(inputTitle) > 185:
        textList = inputTitle.split(' ')
        addNum = len(textList)

        while len(inputTitle) > 180:
            inputTitle = ' '.join(textList[:addNum])
            addNum -= 1

        inputTitle = inputTitle + ' ...'

    # template의 첫번째 텍스트 프레임에 접근
    prs.slides[0].shapes[0].text_frame.paragraphs[0].text = figLabel + inputTitle

    # template의 두번째 텍스트 프레임에 접근
    prs.slides[0].shapes[1].text_frame.paragraphs[0].text = articleInfo

    # 일부에만 서식을 적용하거나 하이퍼링크를 걸고싶을시 사용. 링크된 doi 삽입
    p = prs.slides[0].shapes[1].text_frame.paragraphs[0]
    run = p.add_run()
    run.text = articleDois
    run.hyperlink.address = articleDois

    # 슬라이드 노트에 추가
    # title 첫글자가 볼드나 이탤릭인 경우 \n\n로 시작함. 보정
    if figTitle[:2] == '\n\n':
        prs.slides[0].notes_slide.notes_text_frame.text = figLabel + figTitle[2:].replace('\n\n', '\n')
    else:
        prs.slides[0].notes_slide.notes_text_frame.text = figLabel + figTitle[1:].replace('\n\n', '\n')


# def createthumim(im, filename, filepathImg):
#     percent = 100 / im.width
#     thum_img_resize = im.resize((100, int(im.height * percent)), Image.Resampling.BILINEAR)
#     thum_img_resize.save(filepathImg + '\\Images\\' + filename, dpi=(72, 72))

def createlim(img, l_save_filename, save_dir):
    """
    l-image를 생성하고 지정된 디렉토리에 저장합니다.
    """
    # 저장될 파일의 전체 경로를 os.path.join으로 구성
    save_path = os.path.join(save_dir, l_save_filename)
    factor = min(1, float(150 / img.info['dpi'][0]))
    size = int(factor * img.width), int(factor * img.height)
    l_img_resize = img.resize(size, Image.Resampling.LANCZOS)
    # cios는 해상도를 180으로. 아직 확정 아님
    # if 'cios' in l_save_filename:
    #     l_img_resize.save(save_path, dpi=(180, 180), quality=100)
    # else:
    #     l_img_resize.save(save_path, dpi=(145, 145), quality=100)
    l_img_resize.save(save_path, dpi=(145, 145), quality=100)
def createmim(img, m_save_filename, save_dir):
    """
    m-image를 생성하고 지정된 디렉토리에 저장합니다.
    """
    # 저장될 파일의 전체 경로를 os.path.join으로 구성
    save_path = os.path.join(save_dir, m_save_filename)
    # 가로 비율에 무조건 맞추기
    percent = 270 / img.width

    m_img_resize = img.resize((int(img.width * percent), int(img.height * percent)), Image.Resampling.BILINEAR)
    m_img_resize.save(save_path, dpi=(145, 145), quality=100)